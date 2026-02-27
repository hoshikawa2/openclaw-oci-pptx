from __future__ import annotations

import argparse
import json
import os
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Dict, Any, Optional

from pptx import Presentation


# ============================================================
# PATHS / CONFIG (flexible via env vars and CLI)
# ============================================================

def _env_path(name: str, default: Optional[str] = None) -> Optional[Path]:
    v = os.getenv(name, default)
    return Path(v).expanduser() if v else None


OPENCLAW_WORKDIR = _env_path("OPENCLAW_WORKDIR", ".")  # default: current directory
PPTX_TEMPLATE_PATH = _env_path(
    "PPTX_TEMPLATE_PATH",
    str(OPENCLAW_WORKDIR / "template_openclaw_oci_clean.pptx"),
)
PPTX_OUTPUT_PATH = _env_path(
    "PPTX_OUTPUT_PATH",
    str(OPENCLAW_WORKDIR / "openclaw_oci_presentation.pptx"),
)

# Prefer OCI_CONTENT_FILE (policy name) but accept PPTX_CONTENT_PATH too
PPTX_CONTENT_PATH = _env_path(
    "OCI_CONTENT_FILE",
    os.getenv("PPTX_CONTENT_PATH", str(OPENCLAW_WORKDIR / "content.json")),
)

DEFAULT_LINK = "https://docs.oracle.com/en-us/iaas/Content/generative-ai/home.htm"
DEFAULT_PRESENTER = os.getenv("PPTX_PRESENTER", "Cristiano Hoshikawa")
DEFAULT_COVER_SUBTITLE = os.getenv("PPTX_COVER_SUBTITLE", "Architecture")


# ============================================================
# TEMPLATE ENGINE
# ============================================================

class RedwoodSafePPT:
    """
    Loads a PPTX template, wipes all existing slides safely, and builds a new deck
    using named layouts from the template.
    """

    LAYOUT_COVER = "Cover 1 - Full Image"
    LAYOUT_CONTENT = "Full Page - Light"

    def __init__(self, template_path: Path):
        template_path = Path(template_path).expanduser()
        if not template_path.exists():
            raise FileNotFoundError(f"Template not found: {template_path}")

        self.prs = Presentation(str(template_path))

        # Remove ALL slides
        while len(self.prs.slides) > 0:
            rId = self.prs.slides._sldIdLst[0].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[0]

        self.layouts = {layout.name: layout for layout in self.prs.slide_layouts}

    def _layout(self, name: str):
        if name not in self.layouts:
            available = ", ".join(sorted(self.layouts.keys()))
            raise ValueError(f"Layout '{name}' not found in template. Available: {available}")
        return self.layouts[name]

    def add_content(self, title: str, subhead: str, body: str):

        slide = self.prs.slides.add_slide(self._layout(self.LAYOUT_CONTENT))

        text_placeholders = [
            ph for ph in slide.placeholders
            if ph.has_text_frame
        ]

        if len(text_placeholders) < 2:
            raise RuntimeError("Content layout must have at least 2 text placeholders.")

        text_placeholders[0].text = title
        text_placeholders[1].text = f"{subhead}\n\n{body}"

    def add_cover(self, title, subtitle, presenter):

        slide = self.prs.slides.add_slide(self._layout(self.LAYOUT_COVER))

        text_placeholders = [
            ph for ph in slide.placeholders
            if ph.has_text_frame
        ]

        if len(text_placeholders) < 2:
            raise RuntimeError("Cover layout must have at least 2 text placeholders.")

        text_placeholders[1].text = title
        text_placeholders[3].text = subtitle
        text_placeholders[4].text = presenter
        text_placeholders[2].text = datetime.now().strftime("%d %b %Y")

    def save(self, output_path: Path):
        output_path = Path(output_path).expanduser()
        output_path.parent.mkdir(parents=True, exist_ok=True)
        if output_path.exists():
            output_path.unlink()
        self.prs.save(str(output_path))


# ============================================================
# DECK (fixed 7 slides)
# ============================================================

class OCIStrategicArchitectDeck:
    def __init__(self, template_path: Path):
        self.ppt = RedwoodSafePPT(template_path)

    def _format_section(self, section: Dict[str, Any]) -> str:
        bullets = section.get("bullets", []) or []
        evidence = section.get("evidence", []) or []
        keywords = section.get("keywords", []) or []

        lines = []
        for b in bullets:
            lines.append(f"• {str(b).strip()}")

        # if evidence:
        #     lines.append("")
        #     lines.append("Evidence:")
        #     for e in evidence[:2]:
        #         lines.append(f"- {str(e).strip()}")
        #
        # if keywords:
        #     lines.append("")
        #     lines.append("Keywords: " + ", ".join([str(k).strip() for k in keywords[:8]]))

        return "\n".join(lines).strip()

    def build(self, material_link: str, content: Dict[str, Any], presenter: str, cover_subtitle: str):
        # 1) Cover
        self.ppt.add_cover(
            title=str(content["cover_title"]).strip(),
            subtitle=cover_subtitle,
            presenter=presenter,
        )

        # 2) Intro
        self.ppt.add_content(
            title="Intro",
            subhead="Context and Motivation",
            body=self._format_section(content["introduction"]),
        )

        # 3) Technologies
        self.ppt.add_content(
            title="Technologies",
            subhead="Stack OCI",
            body=self._format_section(content["technologies"]),
        )

        # 4) Architecture
        self.ppt.add_content(
            title="Architecture",
            subhead="Architecture Flow",
            body=self._format_section(content["architecture"]),
        )

        # 5) Problems
        self.ppt.add_content(
            title="Problems",
            subhead="Technical Challenges",
            body=self._format_section(content["problems"]),
        )

        # 6) Demo
        self.ppt.add_content(
            title="Demo",
            subhead="Materials",
            body=f"{material_link}\n\n{self._format_section(content['demo'])}",
        )

        # 7) Conclusion
        self.ppt.add_content(
            title="Conclusion",
            subhead="Strategies",
            body=self._format_section(content["conclusion"]),
        )

        if len(self.ppt.prs.slides) != 7:
            raise RuntimeError("Deck must contain exactly 7 slides.")

    def save(self, output_path: Path):
        self.ppt.save(output_path)


# ============================================================
# CLI
# ============================================================

def _load_json(path: Path) -> Dict[str, Any]:
    path = Path(path).expanduser()
    if not path.exists():
        raise FileNotFoundError(f"content.json not found: {path}")
    with open(path, "r", encoding="utf-8") as f:
        return json.load(f)


def main():
    parser = argparse.ArgumentParser(description="Generate a 7-slide OCI strategic PPTX from a template + content.json.")
    parser.add_argument("--template", default=str(PPTX_TEMPLATE_PATH), help="Path to the PPTX template file.")
    parser.add_argument("--output", default=str(PPTX_OUTPUT_PATH), help="Path to the output PPTX to be written.")
    parser.add_argument("--content", default=str(PPTX_CONTENT_PATH), help="Path to content.json.")
    parser.add_argument("--link", default=os.getenv("OCI_LINK_DEMO", DEFAULT_LINK), help="Source link shown on Demo slide.")
    parser.add_argument("--presenter", default=DEFAULT_PRESENTER, help="Presenter name on cover (if placeholder exists).")
    parser.add_argument("--cover-subtitle", default=DEFAULT_COVER_SUBTITLE, help="Cover subtitle.")
    args = parser.parse_args()

    content = _load_json(Path(args.content))

    deck = OCIStrategicArchitectDeck(Path(args.template))
    deck.build(args.link, content, presenter=args.presenter, cover_subtitle=args.cover_subtitle)
    deck.save(Path(args.output))

    print("✅ PPT generated:", Path(args.output).expanduser().resolve())


if __name__ == "__main__":
    main()
